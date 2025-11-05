"""
PPTX Generation Script

Auto-populates the project review slide deck with results, metrics, and figures
from the trained model evaluation.

Author: FireSentry Team
"""

import json
import pandas as pd
from pathlib import Path
import logging
from datetime import datetime
from typing import Dict, Any, Optional
import sys

# Add project root to path
sys.path.append(str(Path(__file__).parent.parent))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. Install with: pip install python-pptx")

logger = logging.getLogger(__name__)

class PPTXGenerator:
    """
    PowerPoint presentation generator for FireSentry project results.
    
    Auto-populates slides with evaluation results, metrics, and figures.
    """
    
    def __init__(self, 
                 template_path: str = "docs/Mini_Project_review_F1[1].pptx",
                 results_path: str = "docs/results.json"):
        """
        Initialize PPTX generator.
        
        Args:
            template_path: Path to existing PPTX template
            results_path: Path to results JSON file
        """
        self.template_path = Path(template_path)
        self.results_path = Path(results_path)
        self.presentation = None
        self.results = None
        
        logger.info(f"PPTXGenerator initialized with template: {self.template_path}")
    
    def load_results(self):
        """Load evaluation results from JSON file."""
        if not self.results_path.exists():
            logger.error(f"Results file not found: {self.results_path}")
            return False
        
        try:
            with open(self.results_path, 'r') as f:
                self.results = json.load(f)
            logger.info("Results loaded successfully")
            return True
        except Exception as e:
            logger.error(f"Failed to load results: {e}")
            return False
    
    def load_or_create_presentation(self):
        """Load existing presentation or create new one."""
        if self.template_path.exists():
            try:
                self.presentation = Presentation(str(self.template_path))
                logger.info(f"Loaded existing presentation: {self.template_path}")
            except Exception as e:
                logger.warning(f"Failed to load template: {e}. Creating new presentation.")
                self.presentation = Presentation()
        else:
            logger.info("Template not found. Creating new presentation.")
            self.presentation = Presentation()
    
    def add_title_slide(self):
        """Add or update title slide."""
        if len(self.presentation.slides) == 0:
            slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
            slide = self.presentation.slides.add_slide(slide_layout)
        else:
            slide = self.presentation.slides[0]
        
        # Update title
        title = slide.shapes.title
        title.text = "FireSentry: Forest Fire Risk Prediction for Uttarakhand"
        
        # Update subtitle
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = f"MSFS + AutoML Implementation\n{datetime.now().strftime('%B %Y')}"
        
        logger.info("Title slide updated")
    
    def add_dataset_summary_slide(self):
        """Add dataset summary slide."""
        slide_layout = self.presentation.slide_layouts[1]  # Content slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Dataset Summary"
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Dataset information
        dataset_info = [
            "• Geographic Coverage: Uttarakhand, India",
            "• Bounding Box: 28.7°N-31.5°N, 77.5°E-81.0°E",
            "• Temporal Range: 2020-2024 (5 years)",
            "• Data Sources:",
            "  - FIRMS: Active fire detections",
            "  - CHIRPS: Daily precipitation (0.05° resolution)",
            "  - MODIS LST: Land surface temperature",
            "  - MODIS SR: Surface reflectance (NDVI/EVI/NDWI)",
            "  - SRTM: Digital elevation model"
        ]
        
        if self.results:
            model_info = self.results.get('model_info', {})
            dataset_info.extend([
                f"• Total Samples: {self.results.get('performance_metrics', {}).get('accuracy', 'N/A')}",
                f"• Features Used: {model_info.get('features_used', 'N/A')}",
                f"• Model Type: {model_info.get('model_type', 'N/A')}"
            ])
        
        for i, line in enumerate(dataset_info):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
        
        logger.info("Dataset summary slide added")
    
    def add_methodology_slide(self):
        """Add methodology slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Methodology"
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        methodology = [
            "1. Dynamic Time Window (DTW) Algorithm:",
            "   • Identifies drying periods preceding fire events",
            "   • Thresholds: 30mm cumulative, 10mm daily precipitation",
            "   • Maximum lookback: 90 days",
            "",
            "2. Feature Engineering (24 features):",
            "   • Precipitation: min, median, mean, max, sum",
            "   • LST: min, median, mean, max",
            "   • Vegetation Indices: NDVI, EVI, NDWI (min, median, mean, max)",
            "   • Terrain: elevation, slope, aspect",
            "",
            "3. Multi-Stage Feature Selection (MSFS):",
            "   • Stage 1: Mutual Information Gain filtering",
            "   • Stage 2: Recursive Feature Elimination with CV",
            "   • Stage 3: Voting aggregation (n_repeats=3)",
            "",
            "4. Auto-sklearn Ensemble Training:",
            "   • Automated algorithm selection and hyperparameter tuning",
            "   • 5-fold cross-validation",
            "   • Ensemble of multiple algorithms"
        ]
        
        for i, line in enumerate(methodology):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
        
        logger.info("Methodology slide added")
    
    def add_results_slide(self):
        """Add results slide with performance metrics."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Model Performance Results"
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        if self.results:
            metrics = self.results.get('performance_metrics', {})
            
            results_text = [
                "Performance Metrics:",
                f"• Accuracy: {metrics.get('accuracy', 'N/A'):.4f}",
                f"• Precision: {metrics.get('precision', 'N/A'):.4f}",
                f"• Recall: {metrics.get('recall', 'N/A'):.4f}",
                f"• F1-Score: {metrics.get('f1', 'N/A'):.4f}",
                f"• AUC-ROC: {metrics.get('auc', 'N/A'):.4f}",
                "",
                "Feature Selection Results:",
                f"• Original Features: 24",
                f"• Selected Features: {self.results.get('model_info', {}).get('features_used', 'N/A')}",
                f"• Selection Method: MSFS (MIG + RFECV + Voting)",
                "",
                "Selected Features:"
            ]
            
            # Add selected features
            selected_features = self.results.get('model_info', {}).get('selected_features', [])
            for feature in selected_features[:10]:  # Show top 10
                results_text.append(f"  - {feature}")
            
            if len(selected_features) > 10:
                results_text.append(f"  ... and {len(selected_features) - 10} more")
        
        else:
            results_text = ["Results not available. Please run evaluation first."]
        
        for i, line in enumerate(results_text):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
        
        logger.info("Results slide added")
    
    def add_figures_slide(self):
        """Add slide with evaluation figures."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Evaluation Figures"
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        figures_text = [
            "Generated Evaluation Figures:",
            "",
            "1. ROC Curve (fig_roc.png):",
            "   • Shows model's ability to distinguish between fire/no-fire",
            "   • AUC score indicates overall performance",
            "",
            "2. Confusion Matrix (fig_cm.png):",
            "   • True vs Predicted classifications",
            "   • Shows precision and recall breakdown",
            "",
            "3. Feature Importance (fig_feature_importance.png):",
            "   • Most predictive features identified by MSFS",
            "   • Helps understand model decision-making",
            "",
            "4. Precision-Recall Curve (fig_pr_curve.png):",
            "   • Alternative to ROC for imbalanced datasets",
            "   • Shows precision-recall trade-off"
        ]
        
        for i, line in enumerate(figures_text):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
        
        logger.info("Figures slide added")
    
    def add_conclusion_slide(self):
        """Add conclusion and future work slide."""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Conclusion & Future Work"
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        conclusion_text = [
            "Key Achievements:",
            "• Successfully implemented MSFS + AutoML pipeline",
            "• Achieved high prediction accuracy on Uttarakhand fire data",
            "• Built operational FastAPI prediction service",
            "• Validated DTW approach for fire risk assessment",
            "",
            "Future Enhancements:",
            "• Real-time satellite data integration",
            "• Higher spatial resolution (1km grid)",
            "• Anthropogenic factor modeling",
            "• Deep learning architectures (CNN/LSTM)",
            "• Multi-source data fusion",
            "• Climate change scenario projections",
            "• Mobile application for field deployment",
            "• Transfer learning to other Himalayan states"
        ]
        
        for i, line in enumerate(conclusion_text):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
        
        logger.info("Conclusion slide added")
    
    def add_figures_to_slides(self):
        """Add actual figure images to slides if they exist."""
        docs_dir = Path("docs")
        figure_files = {
            "fig_roc.png": "ROC Curve",
            "fig_cm.png": "Confusion Matrix", 
            "fig_feature_importance.png": "Feature Importance",
            "fig_pr_curve.png": "Precision-Recall Curve"
        }
        
        for fig_file, fig_title in figure_files.items():
            fig_path = docs_dir / fig_file
            if fig_path.exists():
                # Create new slide for figure
                slide_layout = self.presentation.slide_layouts[5]  # Blank layout
                slide = self.presentation.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.text = fig_title
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True
                
                # Add image
                try:
                    slide.shapes.add_picture(str(fig_path), Inches(1), Inches(1.5), 
                                           width=Inches(8), height=Inches(5))
                    logger.info(f"Added {fig_file} to presentation")
                except Exception as e:
                    logger.warning(f"Failed to add {fig_file}: {e}")
    
    def generate_presentation(self, output_path: Optional[str] = None):
        """
        Generate complete presentation.
        
        Args:
            output_path: Output file path (default: overwrite template)
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available. Cannot generate presentation.")
            return False
        
        logger.info("Generating PowerPoint presentation")
        
        # Load results
        if not self.load_results():
            logger.warning("Results not available. Creating presentation with placeholder content.")
        
        # Load or create presentation
        self.load_or_create_presentation()
        
        # Add slides
        self.add_title_slide()
        self.add_dataset_summary_slide()
        self.add_methodology_slide()
        self.add_results_slide()
        self.add_figures_slide()
        self.add_conclusion_slide()
        
        # Add figure images
        self.add_figures_to_slides()
        
        # Save presentation
        if output_path is None:
            output_path = self.template_path
        else:
            output_path = Path(output_path)
        
        self.presentation.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")
        
        return True

def main():
    """Main function to generate presentation."""
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not available.")
        print("Install with: pip install python-pptx")
        return
    
    generator = PPTXGenerator()
    
    try:
        success = generator.generate_presentation()
        if success:
            print("PowerPoint presentation generated successfully!")
            print(f"Saved to: {generator.template_path}")
        else:
            print("Failed to generate presentation.")
    except Exception as e:
        logger.error(f"Presentation generation failed: {e}")
        print(f"Error: {e}")

if __name__ == "__main__":
    main()




